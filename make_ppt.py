#!/usr/bin/env python3
"""
make_ppt.py — Enhanced Academic PowerPoint Presentation
Hydrochemical Intelligence Report — Bhubaneswar Groundwater 2024

All values sourced directly from output.log and pipeline CSVs.
Design: Navy/Teal/Gold palette — KIIT academic format.

Authors: Lakshya Nayyar (23053133) & Vaibhav Bhaskar (23053173)
"""

from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─────────────────────────────────────────────────────────────────────
# PATHS
# ─────────────────────────────────────────────────────────────────────
BASE_DIR = Path(__file__).resolve().parent
FIG_DIR  = BASE_DIR / "figures"
OUT_PPT  = BASE_DIR / "Hydrochemical_Analysis_Presentation.pptx"

# ─────────────────────────────────────────────────────────────────────
# COLOUR PALETTE
# ─────────────────────────────────────────────────────────────────────
CLR_NAVY      = RGBColor(0x0D, 0x21, 0x37)   # primary dark navy
CLR_NAVY2     = RGBColor(0x16, 0x37, 0x5A)   # mid navy (panels)
CLR_TEAL      = RGBColor(0x00, 0xB4, 0xD8)   # teal accent
CLR_TEAL_DK   = RGBColor(0x00, 0x87, 0xA6)   # dark teal
CLR_TEAL_LT   = RGBColor(0xD0, 0xF4, 0xFF)   # light teal panel bg
CLR_GOLD      = RGBColor(0xF4, 0xA2, 0x61)   # warm gold highlight
CLR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CLR_OFF_WHITE = RGBColor(0xF4, 0xF7, 0xFC)   # slide background
CLR_DARK      = RGBColor(0x1A, 0x1A, 0x2E)   # near-black body text
CLR_GREY      = RGBColor(0x55, 0x65, 0x75)   # secondary text
CLR_RED       = RGBColor(0xD6, 0x3E, 0x3E)   # danger / non-compliant
CLR_RED_LT    = RGBColor(0xFF, 0xEB, 0xEB)   # light red panel
CLR_GREEN     = RGBColor(0x2A, 0x9D, 0x63)   # safe / good
CLR_GREEN_LT  = RGBColor(0xE8, 0xF8, 0xEE)   # light green panel
CLR_ORANGE    = RGBColor(0xF7, 0x8C, 0x28)   # caution
CLR_PANEL     = RGBColor(0xE8, 0xF2, 0xFC)   # light blue panel bg
CLR_DIVIDER   = RGBColor(0xD0, 0xD8, 0xE8)   # thin divider lines

FONT  = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TITLE_H = Inches(1.2)   # title bar height
CONTENT_TOP = Inches(1.35)

# ─────────────────────────────────────────────────────────────────────
# CORE HELPERS
# ─────────────────────────────────────────────────────────────────────

def safe_pic(slide, img_path, left, top, width=None, height=None):
    p = Path(img_path)
    if not p.exists():
        print(f"  [WARN] Missing figure: {p.name}")
        return None
    kw = {"image_file": str(p), "left": left, "top": top}
    if width:  kw["width"]  = width
    if height: kw["height"] = height
    return slide.shapes.add_picture(**kw)


def set_bg(slide, color: RGBColor):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None,
             rounding=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounding else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def add_text(slide, left, top, width, height, text,
             size=14, bold=False, color=CLR_DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    return tb


def add_bullets(slide, left, top, width, height, items,
                size=14, color=CLR_DARK, bold_first_word=False,
                spacing_after=5, bullet="▶"):
    """Bullet list — minimum 14pt for readability."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(spacing_after)
        p.space_before = Pt(2)
    return tb


def add_caption(slide, left, top, width, text):
    """10pt italic grey caption below a figure."""
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"▲  {text}"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = CLR_GREY
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    return tb


def add_footer(slide, text="Nayyar & Bhaskar  |  KIIT University  |  2025"):
    # Footer strip
    add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), CLR_NAVY)
    # Left: authors
    add_text(slide, Inches(0.4), Inches(7.18), Inches(7), Inches(0.28),
             "Groundwater Quality Analysis — Bhubaneswar 2024",
             size=9, color=CLR_TEAL, align=PP_ALIGN.LEFT)
    # Right: attribution
    add_text(slide, Inches(6.5), Inches(7.18), Inches(6.5), Inches(0.28),
             text, size=9, color=CLR_GREY, align=PP_ALIGN.RIGHT)


def rule(slide, left, top, width=Inches(12.5), color=CLR_TEAL, thickness=Pt(2)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_title_bar(slide, title, subtitle=None, icon=None):
    """Navy gradient title bar."""
    # Main navy bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, TITLE_H, CLR_NAVY)
    # Teal left accent stripe
    add_rect(slide, Inches(0), Inches(0), Inches(0.22), TITLE_H, CLR_TEAL)
    # Gold right accent corner
    add_rect(slide, SLIDE_W - Inches(0.22), Inches(0), Inches(0.22), TITLE_H, CLR_GOLD)
    # Title
    add_text(slide, Inches(0.45), Inches(0.10), Inches(11.5), Inches(0.65),
             title, size=26, bold=True, color=CLR_WHITE)
    # Subtitle
    if subtitle:
        add_text(slide, Inches(0.45), Inches(0.72), Inches(11.5), Inches(0.38),
                 subtitle, size=13, color=CLR_TEAL, italic=True)
    # Bottom rule
    rule(slide, Inches(0), TITLE_H - Pt(3), SLIDE_W, CLR_TEAL, Pt(3))


def stat_box(slide, left, top, width, height, number, label,
             bg=CLR_RED, txt=CLR_WHITE):
    """Coloured callout box: big number + small label."""
    add_rect(slide, left, top, width, height, bg, rounding=True)
    num_h = height * 0.55
    lbl_h = height * 0.40
    add_text(slide, left, top + Pt(4), width, num_h,
             number, size=28, bold=True, color=txt, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + num_h, width, lbl_h,
             label, size=11, bold=True, color=txt, align=PP_ALIGN.CENTER)


def panel_box(slide, left, top, width, height, title, color=CLR_TEAL):
    """Titled content panel."""
    add_rect(slide, left, top, width, Inches(0.38), color)
    add_text(slide, left + Inches(0.12), top + Inches(0.05),
             width - Inches(0.12), Inches(0.28),
             title, size=13, bold=True, color=CLR_WHITE)
    add_rect(slide, left, top + Inches(0.38), width,
             height - Inches(0.38), CLR_PANEL, CLR_DIVIDER)


def interp_box(slide, left, top, width, height, text, bg=CLR_TEAL_LT):
    """Highlighted interpretation box — light teal bg, italic dark text."""
    add_rect(slide, left, top, width, height, bg, CLR_TEAL)
    add_text(slide, left + Inches(0.12), top + Inches(0.08),
             width - Inches(0.24), height - Inches(0.1),
             text, size=12, italic=True, color=CLR_NAVY2,
             align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """SLIDE 1 — Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_NAVY)

    # Decorative teal circle (large, bottom-right corner)
    for r, a in [(Inches(5.5), RGBColor(0x00, 0x5F, 0x73)),
                 (Inches(4.0), RGBColor(0x00, 0x78, 0x99)),
                 (Inches(2.5), RGBColor(0x00, 0x96, 0xBE))]:
        s = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            SLIDE_W - r + Inches(0.5), SLIDE_H - r + Inches(0.5), r, r)
        s.fill.solid(); s.fill.fore_color.rgb = a; s.line.fill.background()

    # Gold top accent stripe
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), CLR_GOLD)

    # University tag — top left
    add_text(slide, Inches(0.6), Inches(0.25), Inches(6), Inches(0.38),
             "KIIT DEEMED TO BE UNIVERSITY  |  SCHOOL OF CIVIL ENGINEERING",
             size=11, bold=True, color=CLR_TEAL, align=PP_ALIGN.LEFT)

    # Academic year — top right
    add_text(slide, Inches(7), Inches(0.25), Inches(6), Inches(0.38),
             "Academic Year  2024 – 25",
             size=11, color=CLR_GREY, align=PP_ALIGN.RIGHT)

    # Horizontal rule
    rule(slide, Inches(1.2), Inches(1.0), Inches(11), CLR_TEAL)

    # Main title
    add_text(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.5),
             "Hydrochemical Intelligence Report",
             size=42, bold=True, color=CLR_WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text(slide, Inches(1), Inches(2.65), Inches(11), Inches(0.75),
             "Multi-Seasonal Groundwater Quality Assessment\nBhubaneswar, Odisha — 2024",
             size=20, color=CLR_TEAL, align=PP_ALIGN.CENTER)

    # Divider
    rule(slide, Inches(3.5), Inches(3.55), Inches(6.3), CLR_GOLD)

    # Parameters strip
    add_text(slide, Inches(0.8), Inches(3.75), Inches(11.5), Inches(0.45),
             "16 Parameters  ·  15 Sampling Locations  ·  3 Seasons  ·  195 Samples  ·  5 ML Models",
             size=13, color=CLR_GREY, align=PP_ALIGN.CENTER)

    # Authors
    add_text(slide, Inches(0.8), Inches(4.35), Inches(11.5), Inches(0.55),
             "Lakshya Nayyar (23053133)     |     Vaibhav Bhaskar (23053173)",
             size=19, bold=True, color=CLR_WHITE, align=PP_ALIGN.CENTER)

    # Guide
    add_text(slide, Inches(0.8), Inches(4.95), Inches(11.5), Inches(0.45),
             "Under the guidance of  Dr. Ajit Kumar Pasayat",
             size=14, color=CLR_GOLD, align=PP_ALIGN.CENTER)

    # Affiliation
    add_text(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4),
             "School of Civil Engineering, KIIT Deemed to be University, Bhubaneswar",
             size=12, color=CLR_GREY, align=PP_ALIGN.CENTER)

    # Bottom rule
    rule(slide, Inches(0), Inches(7.08), SLIDE_W, CLR_GOLD, Pt(4))


def slide_02_outline(prs):
    """SLIDE 2 — Outline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "Presentation Outline",
                  "End-to-End Hydrochemical Intelligence Pipeline — 7 Analytical Tasks")

    items_left = [
        ("01", "Introduction & Study Area",     "15 sites · 3 zones · Bhubaneswar"),
        ("02", "Dataset & Augmentation",        "45 original → 195 samples via CMGP"),
        ("03", "Methodology Pipeline",          "7-task scalable analysis framework"),
        ("04", "Seasonal Dynamics",             "ANOVA/KW tests · % change across seasons"),
        ("05", "IS 10500 Compliance & WQI",     "3-tier BIS assessment · weighted WQI"),
    ]
    items_right = [
        ("06", "Source Analysis",               "Piper · Gibbs · Ionic ratios"),
        ("07", "PCA & Clustering",              "6 PCs (73%) · K-Means k=3"),
        ("08", "ML Forecasting",                "RF · GB · NN · SVR · XGBoost"),
        ("09", "SHAP Explainability",           "Feature importance · GA · NSE · RSR"),
        ("10", "Conclusions & Recommendations", "Policy · Health · Future work"),
    ]

    def draw_items(items, x_start):
        y = CONTENT_TOP
        for num, title, desc in items:
            # Number badge
            s = add_rect(slide, x_start, y, Inches(0.55), Inches(0.72), CLR_TEAL, rounding=True)
            add_text(slide, x_start, y + Inches(0.12), Inches(0.55), Inches(0.45),
                     num, size=18, bold=True, color=CLR_WHITE, align=PP_ALIGN.CENTER)
            # Title
            add_text(slide, x_start + Inches(0.68), y, Inches(5.3), Inches(0.38),
                     title, size=16, bold=True, color=CLR_NAVY)
            # Description
            add_text(slide, x_start + Inches(0.68), y + Inches(0.36),
                     Inches(5.3), Inches(0.3),
                     desc, size=11, color=CLR_GREY, italic=True)
            y += Inches(1.08)

    draw_items(items_left,  Inches(0.45))
    draw_items(items_right, Inches(6.9))

    # Vertical divider
    add_rect(slide, Inches(6.65), CONTENT_TOP, Inches(0.04), Inches(5.4), CLR_DIVIDER)
    add_footer(slide)


def slide_03_introduction(prs):
    """SLIDE 3 — Introduction & Study Area"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "Introduction & Study Area",
                  "Bhubaneswar, Odisha — 15 Groundwater Sampling Locations (2024)")

    # Context paragraph
    interp_box(slide, Inches(0.42), CONTENT_TOP, Inches(12.5), Inches(0.72),
               "Rapid urbanisation, industrial expansion, and unregulated waste disposal in "
               "Bhubaneswar have placed increasing pressure on urban groundwater aquifers. "
               "This study systematically characterises hydrochemical signatures across 3 "
               "land-use zones and 3 seasons, providing a scientific basis for risk management.")

    # Three zone panels
    zones = [
        ("PD — Population Density Areas", CLR_TEAL, CLR_TEAL_LT,
         ["PD-1: Acharya Vihar", "PD-2: Ram Mandir", "PD-3: Sailashree Vihar",
          "PD-4: OMFED Square", "PD-5: Old Town"]),
        ("IA — Industrial Areas", CLR_NAVY2, CLR_PANEL,
         ["IA-1: Mancheswar Industrial Estate", "IA-2: Chandaka Industrial Area",
          "IA-3: OMFED Industries", "IA-4: Rasulgarh", "IA-5: Anmol Industries"]),
        ("DY — Dumping Yards", CLR_ORANGE, RGBColor(0xFF, 0xF3, 0xE3),
         ["DY-1: Bhuasuni Temple Area", "DY-2: Lingaraj Railway Station",
          "DY-3: BMC Micro Compost", "DY-4: Gadakan Road", "DY-5: Daruthenga"]),
    ]

    xp = Inches(0.42)
    for title, hdr_clr, bg_clr, sites in zones:
        panel_box(slide, xp, Inches(2.25), Inches(4.1), Inches(3.0), title, hdr_clr)
        add_rect(slide, xp, Inches(2.25) + Inches(0.38),
                 Inches(4.1), Inches(2.62), bg_clr, CLR_DIVIDER)
        sy = Inches(2.25) + Inches(0.48)
        for site in sites:
            add_text(slide, xp + Inches(0.2), sy, Inches(3.8), Inches(0.38),
                     f"◦  {site}", size=13, color=CLR_DARK)
            sy += Inches(0.48)
        xp += Inches(4.3)

    # Key stats strip
    add_rect(slide, Inches(0.42), Inches(5.45), Inches(12.5), Inches(0.52), CLR_NAVY)
    stats = [
        "📍 15 Locations", "🗓 3 Seasons (Pre/Mon/Post monsoon)", "🔬 16 Parameters",
        "🗃 45 Original + 150 Synthetic = 195 Samples", "⚗ IS 10500:2012 Standard"
    ]
    add_text(slide, Inches(0.55), Inches(5.50), Inches(12.3), Inches(0.44),
             "    |    ".join(stats), size=12, bold=True, color=CLR_TEAL,
             align=PP_ALIGN.CENTER)

    # Parameters row
    add_text(slide, Inches(0.42), Inches(6.08), Inches(12.5), Inches(0.45),
             "Parameters:  pH · EC · TDS · TH · Alkalinity · Ca · Mg · Na · K · Iron · HCO₃ · Cl · SO₄ · NO₃ · F · DO",
             size=13, color=CLR_GREY, align=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_04_dataset(prs):
    """SLIDE 4 — Dataset & Synthetic Augmentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "Dataset Description & Synthetic Augmentation",
                  "CMGP Framework — From 45 Original to 195 Analysis-Ready Samples")

    # Left panel — original
    panel_box(slide, Inches(0.42), CONTENT_TOP, Inches(5.8), Inches(2.65),
              "Original Dataset  (n = 45)", CLR_NAVY2)
    add_bullets(slide, Inches(0.6), CONTENT_TOP + Inches(0.5), Inches(5.5), Inches(2.1),
                ["Source: water quality data_Three Seasons_2024.xlsx",
                 "3 sheets: Premonsoon · Monsoon · Postmonsoon 2024",
                 "15 locations × 3 seasons = 45 samples × 16 parameters",
                 "CBE validity: 93.3%  (mean CBE = 0.76%, range ±10.7%)",
                 "Zero missing values in all chemical parameters",
                 "Cleaned: duplicate headers removed, types enforced"],
                size=13, color=CLR_DARK)

    # Right panel — synthetic
    panel_box(slide, Inches(6.55), CONTENT_TOP, Inches(6.38), Inches(2.65),
              "CMGP Synthetic Augmentation  (n = 150)", CLR_TEAL_DK)
    add_bullets(slide, Inches(6.73), CONTENT_TOP + Inches(0.5), Inches(6.1), Inches(2.1),
                ["Controlled Multivariate Gaussian Perturbation (CMGP)",
                 "Layer 1: Covariance inflation +40%  (wider spread)",
                 "Layer 2: Mean jitter ±6%  (inter-annual variability)",
                 "Layer 3: Multivariate sampling  (preserves correlations)",
                 "Layer 4: Independent Gaussian noise +8%  (measurement error)",
                 "Layer 5: Outlier injection 5% × 2.5σ  (realistic extremes)"],
                size=13, color=CLR_DARK)

    # Stat row
    stats = [
        ("45",    "Original\nSamples",  CLR_NAVY2),
        ("150",   "Synthetic\nSamples", CLR_TEAL_DK),
        ("195",   "Combined\nDataset",  CLR_GREEN),
        ("12.2:1","Sample:\nVariable",  CLR_GOLD),
    ]
    xb = Inches(0.6)
    for num, lbl, clr in stats:
        stat_box(slide, xb, Inches(4.25), Inches(2.85), Inches(0.92),
                 num, lbl, bg=clr)
        xb += Inches(3.15)

    # Figure — original vs synthetic
    safe_pic(slide, BASE_DIR / "fig_original_vs_synthetic_publication.png",
             Inches(0.42), Inches(5.3), width=Inches(8.5), height=Inches(1.9))

    # Defense figure
    safe_pic(slide, FIG_DIR / "task2_validation" / "fig_defense_summary.png",
             Inches(9.1), Inches(5.3), width=Inches(3.85), height=Inches(1.9))
    add_caption(slide, Inches(0.42), Inches(7.06), Inches(8.5),
                "Fig. Original (blue) vs. Synthetic (orange) parameter distributions — CMGP preserves statistical fidelity")

    add_footer(slide)


def slide_05_methodology(prs):
    """SLIDE 5 — Methodology Pipeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "Methodology Pipeline",
                  "7-Task Scalable Analysis Framework — Python · IS 10500:2012 · WHO 2011")

    steps = [
        ("T1", "Data Reconstruction & Cleaning",
         "Fuzzy header mapping · Duplicate/NaN removal · Type coercion · CSV export",
         CLR_NAVY2),
        ("T2", "Data Validation",
         "Missing-value heatmap · IQR + Z-score outlier flags · Charge Balance Error (±10%)",
         CLR_TEAL_DK),
        ("T3", "EDA & Seasonal Dynamics",
         "Descriptive stats · ANOVA / Kruskal-Wallis tests · Seasonal % change · Correlation",
         RGBColor(0x20, 0x7E, 0xA7)),
        ("T4", "IS 10500:2012 Compliance & WQI",
         "3-tier BIS classification · Exceedance factors · WQI (Brown et al., 1970)",
         RGBColor(0x1D, 0x7A, 0x56)),
        ("T5", "Source Analysis",
         "PCA · K-Means (k=3) · Piper diagram · Gibbs mechanism · Ionic ratios",
         RGBColor(0x7B, 0x4F, 0xB2)),
        ("T6", "ML-Based Forecasting",
         "RF · GB · NN · SVR · XGBoost  ·  SHAP explainability  ·  GA / NSE / RSR / MAPE",
         RGBColor(0xBF, 0x5F, 0x00)),
        ("T7", "Scientific Insights",
         "Salinity drivers · Seasonal exceedance ranking · Location risk atlas · Recommendations",
         CLR_RED),
    ]

    y = CONTENT_TOP
    row_h = Inches(0.75)

    for tag, title, desc, clr in steps:
        # Tag badge
        add_rect(slide, Inches(0.42), y, Inches(0.7), row_h, clr, rounding=True)
        add_text(slide, Inches(0.42), y + Inches(0.2), Inches(0.7), Inches(0.4),
                 tag, size=14, bold=True, color=CLR_WHITE, align=PP_ALIGN.CENTER)
        # Title box
        add_rect(slide, Inches(1.25), y, Inches(4.2), row_h, CLR_PANEL, CLR_DIVIDER)
        add_text(slide, Inches(1.38), y + Inches(0.15), Inches(4.0), Inches(0.48),
                 title, size=15, bold=True, color=clr)
        # Description box
        add_rect(slide, Inches(5.6), y, Inches(7.4), row_h, CLR_OFF_WHITE, CLR_DIVIDER)
        add_text(slide, Inches(5.75), y + Inches(0.17), Inches(7.2), Inches(0.42),
                 desc, size=13, color=CLR_DARK)
        # Arrow connector
        add_rect(slide, Inches(0.73), y + row_h, Inches(0.08),
                 Inches(0.04), clr)
        y += row_h + Inches(0.04)

    # Standards bar
    add_rect(slide, Inches(0.42), y + Inches(0.08), Inches(12.5), Inches(0.42), CLR_NAVY)
    add_text(slide, Inches(0.6), y + Inches(0.10), Inches(12.3), Inches(0.35),
             "Standards & References:  IS 10500:2012 (BIS)  ·  WHO Guidelines 4th Ed. (2011)  "
             "·  Brown et al. (1970) WQI  ·  Hair et al. (2010) PCA  ·  Lundberg & Lee (2017) SHAP",
             size=11, bold=True, color=CLR_TEAL, align=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_06a_seasonal(prs):
    """SLIDE 6a — Seasonal Dynamics: Statistics"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "Seasonal Hydrochemical Dynamics — Statistical Analysis",
                  "ANOVA / Kruskal-Wallis Tests · Premonsoon → Monsoon → Postmonsoon")

    # Key finding callouts
    callouts = [
        ("EC +60%", "Premonsoon→Monsoon",  CLR_ORANGE),
        ("TH +497%","Pre→Mon surge",        CLR_RED),
        ("Na +101%","Premonsoon→Monsoon",   CLR_NAVY2),
        ("K +122%", "Pre→Mon (fertiliser)", CLR_TEAL_DK),
    ]
    xb = Inches(0.42)
    for num, lbl, clr in callouts:
        stat_box(slide, xb, CONTENT_TOP, Inches(2.95), Inches(0.88), num, lbl, bg=clr)
        xb += Inches(3.1)

    # Interpretation
    interp_box(slide, Inches(0.42), Inches(2.44), Inches(5.5), Inches(0.9),
               "Monsoon infiltration mobilises surface contaminants into the aquifer. "
               "Na, K, and TH spike sharply due to enhanced rock-water interaction and "
               "potential sewage/fertiliser inputs during the wet season.")

    # Key findings bullets
    add_bullets(slide, Inches(0.42), Inches(3.48), Inches(5.5), Inches(2.6),
                ["All 16 parameters show significant seasonal variation (p < 0.05)",
                 "Monsoon is universally the worst season (210 IS 10500 exceedances)",
                 "EC: Pre=392 → Mon=627 → Post=446 µS/cm",
                 "TDS: Pre=246 → Mon=411 → Post=293 mg/L",
                 "Alkalinity drops −30% Mon→Post (carbonate buffering exhausted)",
                 "DO decreases Premonsoon→Monsoon (−22%) — organic loading"],
                size=14, color=CLR_DARK)

    # Right: boxplot figure
    safe_pic(slide, FIG_DIR / "task3_seasonal" / "fig_seasonal_boxplots.png",
             Inches(6.15), CONTENT_TOP, width=Inches(7.0), height=Inches(3.2))
    add_caption(slide, Inches(6.15), Inches(4.42), Inches(7.0),
                "Fig. Box plots of all 16 parameters by season — Monsoon (orange) consistently highest for EC, TDS, Na, K")

    # Bottom: heatmap + trends
    safe_pic(slide, FIG_DIR / "task3_seasonal" / "fig_seasonal_heatmap.png",
             Inches(0.42), Inches(4.72), width=Inches(6.0), height=Inches(2.3))
    add_caption(slide, Inches(0.42), Inches(7.0), Inches(6.0),
                "Fig. Seasonal mean concentration heatmap — red = high, blue = low")

    safe_pic(slide, FIG_DIR / "task3_seasonal" / "fig_seasonal_trends.png",
             Inches(6.65), Inches(4.72), width=Inches(6.5), height=Inches(2.3))
    add_caption(slide, Inches(6.65), Inches(7.0), Inches(6.5),
                "Fig. Seasonal trends (Pre → Mon → Post) with ±1 SD error bars")

    add_footer(slide)


def slide_06b_seasonal(prs):
    """SLIDE 6b — Distributions, Correlations & Radar"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "Seasonal Variation — Distributions & Correlation Structure",
                  "Violin Plots · Correlation Matrix · Seasonal Radar")

    # Left interpretation
    add_text(slide, Inches(0.42), CONTENT_TOP, Inches(5.2), Inches(0.32),
             "Hydrochemical Interpretation", size=15, bold=True, color=CLR_TEAL)
    add_bullets(slide, Inches(0.42), CONTENT_TOP + Inches(0.35), Inches(5.2), Inches(2.8),
                ["EC ↔ TDS: r = 0.898 — strongest correlation (physical law)",
                 "EC ↔ Na: r = 0.732 — sodium is the dominant conductivity driver",
                 "No other pair exceeds |r| > 0.7 after noise injection",
                 "Violin plots reveal monsoon distributions are right-skewed — "
                   "few samples have very high concentrations",
                 "Premonsoon violin shapes are narrower (less variability, dry season)",
                 "Post-monsoon: transitional chemistry — partial recovery from monsoon peak",
                 "Radar chart: Monsoon polygon area is largest — highest overall chemical load"],
                size=13, color=CLR_DARK)

    # Violin plot
    safe_pic(slide, FIG_DIR / "task3_seasonal" / "fig_seasonal_violins.png",
             Inches(5.85), CONTENT_TOP, width=Inches(7.3), height=Inches(3.0))
    add_caption(slide, Inches(5.85), Inches(4.3), Inches(7.3),
                "Fig. Violin plots — shape width = probability density; wider = more samples at that concentration")

    # Correlation matrix
    safe_pic(slide, FIG_DIR / "task3_seasonal" / "fig_correlation_matrix.png",
             Inches(0.42), Inches(3.65), width=Inches(6.0), height=Inches(3.3))
    add_caption(slide, Inches(0.42), Inches(6.92), Inches(6.0),
                "Fig. Pearson correlation matrix — EC↔TDS (r=0.898) is the strongest pair")

    # Radar chart
    safe_pic(slide, FIG_DIR / "task7_insights" / "fig_seasonal_radar.png",
             Inches(6.65), Inches(4.38), width=Inches(6.5), height=Inches(2.6))
    add_caption(slide, Inches(6.65), Inches(6.92), Inches(6.5),
                "Fig. Seasonal radar — area enclosed = overall chemical load; Monsoon (orange) is largest")

    add_footer(slide)


def slide_07_compliance_wqi(prs):
    """SLIDE 7 — IS 10500 Compliance & WQI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "IS 10500:2012 Compliance & Water Quality Index",
                  "BIS 3-Tier Classification: Safe · Caution · Unsafe  |  WQI (Brown et al., 1970)")

    # Top stat callouts
    callouts = [
        ("75.4%", "pH Unsafe",     CLR_RED),
        ("45.1%", "Iron Unsafe",   CLR_ORANGE),
        ("28.2%", "K Unsafe",      CLR_NAVY2),
        ("11.8%", "NO₃ Unsafe",    RGBColor(0x8B, 0x44, 0x13)),
        ("100%",  "WQI: Good+",    CLR_GREEN),
    ]
    xb = Inches(0.42)
    for num, lbl, clr in callouts:
        stat_box(slide, xb, CONTENT_TOP, Inches(2.38), Inches(0.9), num, lbl, bg=clr)
        xb += Inches(2.5)

    # Compliance interpretation
    interp_box(slide, Inches(0.42), Inches(2.55), Inches(6.0), Inches(0.95),
               "pH non-compliance (75.4%) reflects natural geological acidity of "
               "Odisha's alluvial aquifers — NOT primarily anthropogenic. "
               "Iron (45.1%) is geogenic — laterite soils release Fe²⁺ under reducing conditions. "
               "K (28.2%) excess indicates probable fertiliser (potash) contamination near agricultural zones.")

    # WQI summary
    add_text(slide, Inches(6.65), Inches(2.55), Inches(6.23), Inches(0.32),
             "Water Quality Index — All Samples Good or Better", size=14, bold=True, color=CLR_GREEN)
    add_bullets(slide, Inches(6.65), Inches(2.9), Inches(6.23), Inches(0.98),
                ["Excellent (WQI < 50): 77 samples (39.5%) — Premonsoon dominated",
                 "Good (50–100):         118 samples (60.5%) — Monsoon dominated",
                 "Poor / Very Poor / Unsuitable: 0 samples (0.0%)"],
                size=14, color=CLR_DARK, bullet="✔")

    # 3 figures in a row
    safe_pic(slide, FIG_DIR / "task4_safety" / "fig_is10500_compliance_heatmap.png",
             Inches(0.3),  Inches(3.66), width=Inches(4.15), height=Inches(3.1))
    safe_pic(slide, FIG_DIR / "task4_safety" / "fig_is10500_compliance_bars.png",
             Inches(4.62), Inches(3.66), width=Inches(4.15), height=Inches(3.1))
    safe_pic(slide, FIG_DIR / "task4_safety" / "fig_wqi_analysis.png",
             Inches(8.95), Inches(3.66), width=Inches(4.1),  height=Inches(3.1))

    add_caption(slide, Inches(0.3),  Inches(6.73), Inches(4.15),
                "Fig. Compliance heatmap — season × parameter non-compliance %")
    add_caption(slide, Inches(4.62), Inches(6.73), Inches(4.15),
                "Fig. Stacked bars — Safe (green) / Caution (yellow) / Unsafe (red)")
    add_caption(slide, Inches(8.95), Inches(6.73), Inches(4.1),
                "Fig. WQI histogram + box plot by season + pie chart")

    add_footer(slide)


def slide_08_source(prs):
    """SLIDE 8 — Anthropogenic & Geogenic Source Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "Anthropogenic & Geogenic Source Analysis",
                  "Dominant Facies: Ca-Cl Type  |  Mechanism: Rock-Water Interaction")

    # Left bullets
    add_text(slide, Inches(0.42), CONTENT_TOP, Inches(5.5), Inches(0.32),
             "Source Attribution Findings", size=15, bold=True, color=CLR_TEAL)
    add_bullets(slide, Inches(0.42), CONTENT_TOP + Inches(0.35), Inches(5.5), Inches(3.15),
                ["Dominant facies: Ca-Cl type (Piper diamond — bottom-right zone)",
                 "Gibbs mechanism: Rock-water interaction — neither evaporation nor precipitation",
                 "Mean Na/(Na+Ca) = 0.444 < 0.5 → confirms weathering control",
                 "Mean Cl/(Cl+HCO₃) = 0.433 → supports weathering over evaporation",
                 "Ca/Mg mean = 3.28 → calcite dissolution dominates over dolomite",
                 "Na/Cl mean = 0.89 (<1) → Cl excess suggests anthropogenic input or evaporite",
                 "NO₃: mean = 23.8 mg/L, max = 95.2 mg/L → agricultural / sewage signal",
                 "K enrichment (28.2% non-compliant) → probable fertiliser (potash) origin",
                 "Na-Cl excess = −0.60 meq/L → anthropogenic Cl pathway identified"],
                size=13, color=CLR_DARK)

    # Piper diagram (large)
    safe_pic(slide, FIG_DIR / "task5_source" / "fig_piper_diagram.png",
             Inches(6.15), CONTENT_TOP, width=Inches(3.55), height=Inches(3.1))
    add_caption(slide, Inches(6.15), Inches(4.45), Inches(3.55),
                "Fig. Piper diagram — Ca-Cl facies (bottom-right diamond)")

    # Gibbs diagram
    safe_pic(slide, FIG_DIR / "task5_source" / "fig_gibbs_diagram.png",
             Inches(9.85), CONTENT_TOP, width=Inches(3.3), height=Inches(3.1))
    add_caption(slide, Inches(9.85), Inches(4.45), Inches(3.3),
                "Fig. Gibbs diagram — samples in rock-water interaction zone")

    # Ionic ratios (bottom left)
    safe_pic(slide, FIG_DIR / "task5_source" / "fig_ionic_ratios.png",
             Inches(0.3),  Inches(4.62), width=Inches(5.8), height=Inches(2.5))
    add_caption(slide, Inches(0.3), Inches(7.08), Inches(5.8),
                "Fig. Na/Cl and Ca/Mg ionic ratios by season — Ca/Mg >2 confirms calcite dissolution")

    # Ionic scatter
    safe_pic(slide, FIG_DIR / "task5_source" / "fig_piper_ternary.png",
             Inches(6.3),  Inches(4.62), width=Inches(6.85), height=Inches(2.5))
    add_caption(slide, Inches(6.3), Inches(7.08), Inches(6.85),
                "Fig. Piper ternary subplots — cation and anion percentage composition per season")

    add_footer(slide)


def slide_09_pca_clustering(prs):
    """SLIDE 9 — PCA & K-Means Clustering"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "Parameter Grouping — PCA & K-Means Clustering",
                  "6 Principal Components Explain 73.15% of Total Variance  |  3 Hydrochemical Clusters")

    # PCA summary
    add_text(slide, Inches(0.42), CONTENT_TOP, Inches(6.0), Inches(0.32),
             "Principal Component Analysis", size=15, bold=True, color=CLR_TEAL)
    add_bullets(slide, Inches(0.42), CONTENT_TOP + Inches(0.35), Inches(6.0), Inches(2.0),
                ["PC1 (29.3%): Overall mineralisation — high loadings for EC, TDS, Ca, Na",
                 "PC2 (13.6%): Carbonate buffering — Alkalinity, HCO₃ dominant",
                 "PC3 (9.8%):  Geogenic Mg vs. anthropogenic NO₃ (negative loading)",
                 "PC4 (8.2%):  Redox conditions — Iron loading = 0.629",
                 "6 PCs retained (>73% variance) — scree elbow confirms selection"],
                size=14, color=CLR_DARK)

    # Clustering summary
    add_text(slide, Inches(6.65), CONTENT_TOP, Inches(6.38), Inches(0.32),
             "K-Means Clustering (k = 3)", size=15, bold=True, color=CLR_TEAL)
    add_bullets(slide, Inches(6.65), CONTENT_TOP + Inches(0.35), Inches(6.38), Inches(2.0),
                ["Elbow method confirms k=3 as optimal — confirmed by Ward dendrogram",
                 "Cluster 0 — High mineralisation: TDS ≈ 403, EC ≈ 620 µS/cm",
                 "Cluster 1 — Low mineralisation:  TDS ≈ 251, EC ≈ 390 µS/cm",
                 "Cluster 2 — Moderate / hardness-enriched: TDS ≈ 356, EC ≈ 559",
                 "Monsoon samples dominate Cluster 0 — seasonal enrichment confirmed"],
                size=14, color=CLR_DARK)

    # Figures row 1
    safe_pic(slide, FIG_DIR / "task5_source" / "fig_pca_scree.png",
             Inches(0.3),  Inches(3.5), width=Inches(3.2), height=Inches(2.5))
    safe_pic(slide, FIG_DIR / "task5_source" / "fig_pca_biplot.png",
             Inches(3.65), Inches(3.5), width=Inches(3.2), height=Inches(2.5))
    safe_pic(slide, FIG_DIR / "task5_source" / "fig_kmeans_pca.png",
             Inches(7.0),  Inches(3.5), width=Inches(3.2), height=Inches(2.5))
    safe_pic(slide, FIG_DIR / "task5_source" / "fig_dendrogram.png",
             Inches(10.3), Inches(3.5), width=Inches(2.8), height=Inches(2.5))

    add_caption(slide, Inches(0.3),  Inches(5.95), Inches(3.2),
                "Fig. Scree plot — elbow at PC6 (73.15% variance)")
    add_caption(slide, Inches(3.65), Inches(5.95), Inches(3.2),
                "Fig. PCA biplot — dots=samples, arrows=parameters")
    add_caption(slide, Inches(7.0),  Inches(5.95), Inches(3.2),
                "Fig. K-Means clusters on PC1-PC2 space")
    add_caption(slide, Inches(10.3), Inches(5.95), Inches(2.8),
                "Fig. Ward dendrogram — confirms k=3")

    # PCA loadings heatmap
    safe_pic(slide, FIG_DIR / "task5_source" / "fig_pca_loadings_heatmap.png",
             Inches(0.3), Inches(6.3), width=Inches(8.5), height=Inches(0.78))
    add_caption(slide, Inches(0.3), Inches(7.05), Inches(8.5),
                "Fig. PCA loading heatmap — red=positive/dominant, blue=negative contribution per component")

    safe_pic(slide, FIG_DIR / "task5_source" / "fig_elbow.png",
             Inches(9.0), Inches(6.3), width=Inches(4.1), height=Inches(0.78))

    add_footer(slide)


def slide_10a_ml(prs):
    """SLIDE 10a — ML Model Comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "ML-Based Forecasting — Model Comparison",
                  "5 Models × 3 Targets (TDS, EC, WQI)  |  80:20 Split  |  5-Fold CV")

    # Best model callouts
    callouts = [
        ("RF",   "Best: TDS\nCV R²=0.753",  CLR_NAVY2),
        ("RF",   "Best: EC\nCV R²=0.771",   CLR_TEAL_DK),
        ("NN",   "Best: WQI\nCV R²=0.903",  CLR_GREEN),
        ("0.764","NSE (TDS-RF)",             CLR_ORANGE),
        ("8.8%", "MAPE (EC-RF)",             RGBColor(0x6A, 0x3D, 0x9A)),
    ]
    xb = Inches(0.42)
    for num, lbl, clr in callouts:
        stat_box(slide, xb, CONTENT_TOP, Inches(2.38), Inches(0.9), num, lbl, bg=clr)
        xb += Inches(2.5)

    # Per-target result columns
    targets = [
        ("TARGET: TDS", [
            "Random Forest:  CV R²=0.753  ·  Test R²=0.764  ·  RMSE=49.3  ·  MAPE=13.1%",
            "Gradient Boost: CV R²=0.674  ·  Test R²=0.764  ·  RMSE=49.2",
            "XGBoost:        CV R²=0.714  ·  Test R²=0.764  ·  RMSE=49.3",
            "Neural Network: CV R²=0.704  ·  Test R²=0.737  ·  RMSE=52.0",
            "SVR:            CV R²=0.711  ·  Test R²=0.709  ·  RMSE=54.7",
        ]),
        ("TARGET: EC", [
            "Random Forest:  CV R²=0.771  ·  Test R²=0.757  ·  RMSE=67.6  ·  MAPE=8.8%",
            "Gradient Boost: CV R²=0.747  ·  Test R²=0.751  ·  RMSE=68.5",
            "XGBoost:        CV R²=0.742  ·  Test R²=0.718  ·  RMSE=72.9",
            "Neural Network: CV R²=0.740  ·  Test R²=0.725  ·  RMSE=72.0",
            "SVR:            CV R²=0.715  ·  Test R²=0.736  ·  RMSE=70.5",
        ]),
        ("TARGET: WQI", [
            "Neural Network: CV R²=0.903  ·  Test R²=0.954  ·  RMSE=2.1   ·  MAPE=3.2%",
            "SVR:            CV R²=0.883  ·  Test R²=0.951  ·  RMSE=2.1",
            "XGBoost:        CV R²=0.830  ·  Test R²=0.886  ·  RMSE=3.3",
            "Gradient Boost: CV R²=0.830  ·  Test R²=0.881  ·  RMSE=3.3",
            "Random Forest:  CV R²=0.800  ·  Test R²=0.863  ·  RMSE=3.6",
        ]),
    ]
    xp = Inches(0.3)
    for title, items in targets:
        add_text(slide, xp, Inches(2.55), Inches(4.3), Inches(0.34),
                 title, size=14, bold=True, color=CLR_NAVY2)
        add_rect(slide, xp, Inches(2.9), Inches(4.3), Inches(1.75), CLR_PANEL, CLR_DIVIDER)
        iy = Inches(2.95)
        for line in items:
            clr = CLR_TEAL if "Random Forest" in line or "Neural Network" in line else CLR_DARK
            bold = "Random Forest" in line or "Neural Network" in line
            add_text(slide, xp + Inches(0.1), iy, Inches(4.1), Inches(0.3),
                     line, size=11, bold=bold, color=clr)
            iy += Inches(0.33)
        xp += Inches(4.42)

    # Actual vs predicted figures
    safe_pic(slide, FIG_DIR / "task6_ml" / "fig_actual_vs_predicted_tds.png",
             Inches(0.3),  Inches(4.82), width=Inches(4.2), height=Inches(2.2))
    safe_pic(slide, FIG_DIR / "task6_ml" / "fig_actual_vs_predicted_ec.png",
             Inches(4.6),  Inches(4.82), width=Inches(4.2), height=Inches(2.2))
    safe_pic(slide, FIG_DIR / "task6_ml" / "fig_actual_vs_predicted_wqi.png",
             Inches(9.0),  Inches(4.82), width=Inches(4.1), height=Inches(2.2))

    add_caption(slide, Inches(0.3),  Inches(6.98), Inches(4.2),
                "Fig. Actual vs. Predicted TDS — points on 1:1 line = perfect prediction")
    add_caption(slide, Inches(4.6),  Inches(6.98), Inches(4.2),
                "Fig. Actual vs. Predicted EC — RF (blue) hugs the line most tightly")
    add_caption(slide, Inches(9.0),  Inches(6.98), Inches(4.1),
                "Fig. Actual vs. Predicted WQI — NN achieves R²=0.954")

    add_footer(slide)


def slide_10b_shap(prs):
    """SLIDE 10b — SHAP Explainability & Residuals"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "SHAP Explainability, Residuals & Uncertainty Metrics",
                  "SHapley Additive exPlanations — Interpreting Black-Box ML Predictions")

    # SHAP explanation
    interp_box(slide, Inches(0.42), CONTENT_TOP, Inches(12.5), Inches(0.72),
               "SHAP decomposes each model prediction into per-feature contributions. "
               "Each dot = one sample; X-position = how much that feature pushed "
               "the prediction up (right) or down (left); colour = feature value (red=high, blue=low). "
               "Ca and Na consistently push TDS/EC upward across all samples — confirming Ca-Cl geochemistry.")

    # Feature importance bullets
    add_text(slide, Inches(0.42), Inches(2.38), Inches(5.5), Inches(0.32),
             "Top SHAP Features", size=15, bold=True, color=CLR_TEAL)
    add_bullets(slide, Inches(0.42), Inches(2.73), Inches(5.5), Inches(1.65),
                ["TDS:  Ca (42.5%) > Na (28.6%) > K (6.8%)",
                 "EC:   Na (45.9%) > Ca (21.0%) > K (11.8%)",
                 "WQI:  EC (40.0%) > TDS (20.9%) > NO₃ (13.2%)",
                 "Uncertainty (TDS-RF):  GA=0.788 · NSE=0.764 · RSR=0.486 · MAPE=13.1%",
                 "Uncertainty (EC-RF):   GA=0.781 · NSE=0.757 · RSR=0.493 · MAPE=8.8%"],
                size=14, color=CLR_DARK)

    # Feature importance figure
    safe_pic(slide, FIG_DIR / "task6_ml" / "fig_feature_importance.png",
             Inches(6.0), Inches(2.35), width=Inches(7.15), height=Inches(2.1))
    add_caption(slide, Inches(6.0), Inches(4.4), Inches(7.15),
                "Fig. Gini-based feature importance (RF) for TDS, EC, WQI — Ca and Na dominate")

    # SHAP summary plots (beeswarm)
    safe_pic(slide, FIG_DIR / "task6_ml" / "fig_shap_summary_tds.png",
             Inches(0.3),  Inches(4.55), width=Inches(4.25), height=Inches(1.9))
    safe_pic(slide, FIG_DIR / "task6_ml" / "fig_shap_summary_ec.png",
             Inches(4.65), Inches(4.55), width=Inches(4.25), height=Inches(1.9))
    safe_pic(slide, FIG_DIR / "task6_ml" / "fig_shap_summary_wqi.png",
             Inches(9.05), Inches(4.55), width=Inches(4.1),  height=Inches(1.9))

    add_caption(slide, Inches(0.3),  Inches(6.40), Inches(4.25),
                "Fig. SHAP beeswarm (TDS) — red=high Ca/Na pushes TDS up")
    add_caption(slide, Inches(4.65), Inches(6.40), Inches(4.25),
                "Fig. SHAP beeswarm (EC) — Na dominates EC predictions")
    add_caption(slide, Inches(9.05), Inches(6.40), Inches(4.1),
                "Fig. SHAP beeswarm (WQI) — EC and TDS are key WQI drivers")

    # Residuals strip
    safe_pic(slide, FIG_DIR / "task6_ml" / "fig_residuals_tds.png",
             Inches(0.3),  Inches(6.62), width=Inches(4.25), height=Inches(0.62))
    safe_pic(slide, FIG_DIR / "task6_ml" / "fig_residuals_ec.png",
             Inches(4.65), Inches(6.62), width=Inches(4.25), height=Inches(0.62))
    safe_pic(slide, FIG_DIR / "task6_ml" / "fig_residuals_wqi.png",
             Inches(9.05), Inches(6.62), width=Inches(4.1),  height=Inches(0.62))

    add_footer(slide)


def slide_11_conclusion(prs):
    """SLIDE 11 — Conclusions & Recommendations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "Conclusions & Policy Recommendations",
                  "Evidence-Based Findings from 7-Task Hydrochemical Intelligence Pipeline")

    # Key findings
    add_text(slide, Inches(0.42), CONTENT_TOP, Inches(6.0), Inches(0.34),
             "Key Scientific Findings", size=16, bold=True, color=CLR_TEAL)
    add_bullets(slide, Inches(0.42), CONTENT_TOP + Inches(0.38), Inches(6.0), Inches(4.42),
                ["pH (75.4% unsafe) — natural geological acidity of Odisha's alluvial aquifers",
                 "Iron (45.1% unsafe) — geogenic, laterite soils release Fe²⁺ under reducing conditions",
                 "K (28.2%) and NO₃ (11.8%) non-compliance — probable agricultural (fertiliser) inputs",
                 "Dominant facies: Ca-Cl type — rock-water interaction controls chemistry (Gibbs zone)",
                 "Monsoon is worst season — 210 IS 10500 exceedances vs 88 in Premonsoon",
                 "WQI range 30.81–77.04 — all 195 samples in 'Excellent' or 'Good' categories",
                 "Best ML models: RF (TDS CV R²=0.753), RF (EC CV R²=0.771), NN (WQI CV R²=0.903)",
                 "SHAP confirms Ca and Na as dominant TDS/EC drivers — consistent with Ca-Cl facies",
                 "IA-3 is most contaminated site: TDS=393.6 mg/L, EC=607.5 µS/cm, pH=5.87"],
                size=13, color=CLR_DARK, spacing_after=4)

    # Policy recommendations
    add_text(slide, Inches(6.65), CONTENT_TOP, Inches(6.38), Inches(0.34),
             "Policy Recommendations", size=16, bold=True, color=CLR_GOLD)
    rec_clr = RGBColor(0x4A, 0x2C, 0x01)
    add_bullets(slide, Inches(6.65), CONTENT_TOP + Inches(0.38), Inches(6.3), Inches(2.5),
                ["pH correction: lime/soda-ash alkaline dosing at extraction points",
                 "Regulate K-based fertiliser use near IA zone groundwater catchments",
                 "Establish real-time monitoring at IA-3, DY-4, IA-1 (highest risk)",
                 "Rainwater harvesting to boost natural aquifer recharge rates",
                 "Manage organic waste at DY sites to prevent acidification",
                 "Deploy RF model for TDS/EC early-warning forecasting",
                 "Climate-adaptive seasonal monitoring — monsoon variability critical"],
                size=13, color=CLR_DARK)

    # Health implications panel
    add_rect(slide, Inches(6.65), Inches(4.1), Inches(6.38), Inches(2.7), CLR_RED_LT, CLR_RED)
    add_text(slide, Inches(6.78), Inches(4.15), Inches(6.1), Inches(0.35),
             "⚠ Health Implications", size=14, bold=True, color=CLR_RED)
    add_bullets(slide, Inches(6.78), Inches(4.5), Inches(6.1), Inches(2.25),
                ["pH < 6.5 — corrosive water, metal leaching risk (Cu, Pb from pipes)",
                 "NO₃ > 45 mg/L — methemoglobinemia (blue-baby syndrome) risk",
                 "Iron > 0.3 mg/L — taste issues, promotes bacteria growth (51.8% samples)",
                 "Low DO — indicates organic contamination load in monsoon"],
                size=13, color=CLR_DARK)

    add_footer(slide)


def slide_12_acknowledgement(prs):
    """SLIDE 12 — Acknowledgement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_NAVY)

    # Decorative circle
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        SLIDE_W - Inches(3.5), SLIDE_H - Inches(3.5), Inches(4.5), Inches(4.5))
    s.fill.solid(); s.fill.fore_color.rgb = CLR_NAVY2; s.line.fill.background()

    # Gold top rule
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.14), CLR_GOLD)

    add_text(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.75),
             "Acknowledgement",
             size=38, bold=True, color=CLR_WHITE, align=PP_ALIGN.CENTER)

    rule(slide, Inches(4), Inches(1.2), Inches(5), CLR_TEAL)

    ack = (
        "We extend our sincere gratitude to\n\n"
        "Dr. Ajit Kumar Pasayat\n"
        "for his invaluable guidance, mentorship, and continuous support\n"
        "throughout the course of this project.\n\n"
        "We also thank the School of Civil Engineering,\n"
        "KIIT Deemed to be University, Bhubaneswar,\n"
        "for providing laboratory facilities and field sampling infrastructure.\n\n"
        "This study was conducted as part of the Academic Year 2024–25 coursework."
    )
    add_text(slide, Inches(1.5), Inches(1.55), Inches(10), Inches(4.0),
             ack, size=17, color=CLR_WHITE, align=PP_ALIGN.CENTER)

    rule(slide, Inches(3), Inches(5.78), Inches(7.3), CLR_GOLD)

    add_text(slide, Inches(1), Inches(5.98), Inches(11), Inches(0.5),
             "Lakshya Nayyar (23053133)     |     Vaibhav Bhaskar (23053173)",
             size=16, bold=True, color=CLR_TEAL, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(6.52), Inches(11), Inches(0.45),
             "School of Civil Engineering  ·  KIIT Deemed to be University  ·  Bhubaneswar",
             size=12, color=CLR_GREY, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0), Inches(7.08), SLIDE_W, Inches(0.42), CLR_GOLD)
    add_text(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.38),
             "Thank You", size=18, bold=True, color=CLR_NAVY, align=PP_ALIGN.CENTER)


def slide_13_references(prs):
    """SLIDE 13 — References"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CLR_OFF_WHITE)
    add_title_bar(slide, "References",
                  "Standards · Methodological Sources · Literature")

    refs = [
        ("IS 10500:2012",  "Bureau of Indian Standards. Drinking Water Specification, Second Revision. New Delhi: BIS, 2012."),
        ("WHO 2011",       "World Health Organization. Guidelines for Drinking-water Quality, 4th Edition. Geneva: WHO, 2011."),
        ("Brown 1970",     "Brown R.M. et al. A Water Quality Index — Do We Dare? Water & Sewage Works, 117(10), 1970."),
        ("Sekar 2025",     "Sekar S. et al. ML-based prediction of seasonal groundwater quality for Melur, Tamil Nadu. Results in Engineering, 28, 2025."),
        ("Piper 1944",     "Piper A.M. A Graphic Procedure in the Geochemical Interpretation of Water Analyses. Trans. AGU, 25(6), 914–928, 1944."),
        ("Gibbs 1970",     "Gibbs R.J. Mechanisms Controlling World Water Chemistry. Science, 170(3962), 1088–1090, 1970."),
        ("Lundberg 2017",  "Lundberg S.M. & Lee S.-I. A Unified Approach to Interpreting Model Predictions. NeurIPS 2017."),
        ("Hair 2010",      "Hair J.F. et al. Multivariate Data Analysis. Pearson, 7th Edition, 2010."),
        ("CPCB / FSSAI",   "Central Pollution Control Board. Drinking Water Quality — Class A Criteria. Govt. of India."),
    ]

    y = CONTENT_TOP
    for tag, text in refs:
        add_rect(slide, Inches(0.42), y, Inches(1.8), Inches(0.45), CLR_TEAL, rounding=True)
        add_text(slide, Inches(0.42), y + Inches(0.06), Inches(1.8), Inches(0.35),
                 tag, size=11, bold=True, color=CLR_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(2.35), y + Inches(0.04), Inches(10.6), Inches(0.4),
                 text, size=12, color=CLR_DARK)
        y += Inches(0.58)

    add_footer(slide)


# ─────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────

def main():
    print("=" * 70)
    print("  GENERATING ENHANCED POWERPOINT PRESENTATION")
    print("=" * 70)

    prs = Presentation()
    prs.slide_width  = Emu(12192000)   # 13.333 inches
    prs.slide_height = Emu(6858000)    # 7.5 inches

    builders = [
        ("Slide 01: Title",                  slide_01_title),
        ("Slide 02: Outline",                slide_02_outline),
        ("Slide 03: Introduction",           slide_03_introduction),
        ("Slide 04: Dataset & Augmentation", slide_04_dataset),
        ("Slide 05: Methodology",            slide_05_methodology),
        ("Slide 06a: Seasonal — Stats",      slide_06a_seasonal),
        ("Slide 06b: Seasonal — Corr/Radar", slide_06b_seasonal),
        ("Slide 07: IS 10500 & WQI",         slide_07_compliance_wqi),
        ("Slide 08: Source Analysis",        slide_08_source),
        ("Slide 09: PCA & Clustering",       slide_09_pca_clustering),
        ("Slide 10a: ML Models",             slide_10a_ml),
        ("Slide 10b: SHAP & Residuals",      slide_10b_shap),
        ("Slide 11: Conclusions",            slide_11_conclusion),
        ("Slide 12: Acknowledgement",         slide_12_acknowledgement),
        ("Slide 13: References",             slide_13_references),
    ]

    for label, builder in builders:
        print(f"  Building {label} ...")
        builder(prs)

    prs.save(str(OUT_PPT))
    size_kb = OUT_PPT.stat().st_size / 1024
    print(f"\n  ✓ Saved: {OUT_PPT}")
    print(f"  ✓ Size:  {size_kb:.1f} KB")
    print(f"  ✓ Slides: {len(prs.slides)}")
    print("=" * 70)


if __name__ == "__main__":
    main()

