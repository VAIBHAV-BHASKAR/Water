"""
Generate Word Report + PowerPoint Presentation
Hydrochemical Analysis — Bhubaneswar, Odisha, India (2024)
"""
import os, pathlib
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = pathlib.Path(__file__).parent
FIG = BASE / "figures"

# ── helpers ──────────────────────────────────────────────────────────────
def _add_fig_docx(doc, path, caption, width=5.5):
    if path.exists():
        doc.add_picture(str(path), width=Inches(width))
        last = doc.paragraphs[-1]
        last.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p = doc.add_paragraph(caption)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.runs[0].font.size = Pt(9)
        p.runs[0].font.italic = True

def _add_table_docx(doc, headers, rows):
    t = doc.add_table(rows=1+len(rows), cols=len(headers))
    t.style = 'Light Shading Accent 1'
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        c = t.rows[0].cells[i]
        c.text = h
        for r in c.paragraphs:
            for run in r.runs:
                run.font.bold = True
                run.font.size = Pt(9)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = t.rows[ri+1].cells[ci]
            cell.text = str(val)
            for r in cell.paragraphs:
                for run in r.runs:
                    run.font.size = Pt(9)

def _slide(prs, title, bullets=None, fig_path=None, layout_idx=1):
    """Add a slide. layout 0=title, 1=title+content, 5=blank, 6=title only."""
    slide = prs.slides.add_slide(prs.slide_layouts[min(layout_idx, len(prs.slide_layouts)-1)])
    # title
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        tf.text = title
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = PPt(28)
                r.font.bold = True
                r.font.color.rgb = PRGBColor(0, 51, 102)
    # bullets
    if bullets and layout_idx == 1:
        body = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body = ph.text_frame
                break
        if body:
            body.clear()
            for i, b in enumerate(bullets):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = b
                p.font.size = PPt(16)
                p.space_after = PPt(4)
    # figure
    if fig_path and fig_path.exists():
        left = PInches(0.5)
        if bullets:
            top = PInches(5.0)
            h = PInches(2.3)
        else:
            top = PInches(1.6)
            h = PInches(5.2)
        slide.shapes.add_picture(str(fig_path), left, top, height=h)
    return slide

def _slide_fig_only(prs, title, fig_path, caption=""):
    """Full-slide figure with small caption."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # title box
    from pptx.util import Inches as I2
    txBox = slide.shapes.add_textbox(I2(0.3), I2(0.15), I2(9.4), I2(0.6))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = PPt(22); p.font.bold = True; p.font.color.rgb = PRGBColor(0,51,102)
    # figure
    if fig_path and fig_path.exists():
        slide.shapes.add_picture(str(fig_path), I2(0.4), I2(0.85), height=I2(5.8))
    # caption
    if caption:
        cb = slide.shapes.add_textbox(I2(0.3), I2(6.8), I2(9.4), I2(0.5))
        cp = cb.text_frame.paragraphs[0]; cp.text = caption
        cp.font.size = PPt(11); cp.font.italic = True
    return slide

# ═══════════════════════════════════════════════════════════════════════
# WORD REPORT
# ═══════════════════════════════════════════════════════════════════════
def build_word():
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)

    # ── Title page ──
    for _ in range(6):
        doc.add_paragraph()
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("Hydrochemical Intelligence Pipeline\nfor Groundwater Quality Analysis")
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = RGBColor(0,51,102)
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = sub.add_run("Bhubaneswar, Odisha, India (2024)")
    r2.font.size = Pt(16); r2.font.color.rgb = RGBColor(80,80,80)
    sub2 = doc.add_paragraph()
    sub2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r3 = sub2.add_run("15 Locations · 16 Parameters · 3 Seasons")
    r3.font.size = Pt(12); r3.font.italic = True
    doc.add_page_break()

    # ── 1. INTRODUCTION ──
    doc.add_heading("1. Introduction", level=1)
    doc.add_paragraph(
        "Groundwater is the primary drinking water source for over 85% of India's rural population. "
        "In Bhubaneswar, the capital of Odisha, rapid urbanization, industrial expansion, and "
        "agricultural intensification have placed increasing pressure on groundwater resources. "
        "Hydrochemical characterization is essential to assess water potability, identify contamination "
        "sources (geogenic vs. anthropogenic), and formulate evidence-based management strategies."
    )
    doc.add_paragraph(
        "This study presents a comprehensive, reproducible hydrochemical intelligence pipeline covering "
        "data cleaning, statistical validation, seasonal trend analysis, regulatory compliance assessment "
        "(IS 10500:2012), source identification (geogenic vs. anthropogenic), and machine learning–based "
        "water quality prediction. The pipeline processes data from 15 groundwater monitoring locations "
        "across three land-use categories — Population Density (PD), Industrial Area (IA), and Dumping "
        "Yard (DY) — sampled during Premonsoon, Monsoon, and Postmonsoon seasons of 2024."
    )
    doc.add_heading("1.1 Objectives", level=2)
    objectives = [
        "Clean and validate raw hydrochemical data from field sampling campaigns.",
        "Augment the dataset using Conditional Multivariate Gaussian Perturbation (CMGP) for statistical robustness.",
        "Characterize seasonal variations across all 16 hydrochemical parameters.",
        "Assess compliance with IS 10500:2012 Drinking Water Standards and compute Water Quality Index (WQI).",
        "Identify geogenic vs. anthropogenic contamination sources using multi-evidence convergence.",
        "Develop machine learning models (RF, GB, NN, SVR, XGBoost) for TDS, EC, and WQI prediction.",
        "Provide actionable remediation recommendations for Bhubaneswar's groundwater management.",
    ]
    for o in objectives:
        doc.add_paragraph(o, style='List Bullet')

    doc.add_heading("1.2 Study Area", level=2)
    doc.add_paragraph(
        "Bhubaneswar (20.27°N, 85.84°E) lies in the Eastern Ghats mobile belt, underlain by "
        "charnockite-khondalite group rocks with a lateritic weathering cap. The aquifer system "
        "consists of weathered and fractured crystalline rocks overlain by alluvial deposits. "
        "The lateritic profile contributes naturally elevated iron and acidic pH to groundwater. "
        "The study covers 15 locations across three area types: PD-1 to PD-5 (Population Density), "
        "IA-1 to IA-5 (Industrial Area), and DY-1 to DY-5 (Dumping Yard)."
    )
    doc.add_page_break()

    # ── 2. DATA ──
    doc.add_heading("2. Data Description", level=1)
    doc.add_paragraph(
        "Raw hydrochemical data were collected from 15 monitoring locations across 3 seasons "
        "(Premonsoon, Monsoon, Postmonsoon), yielding 45 original samples. Each sample was "
        "analyzed for 16 physicochemical parameters."
    )
    _add_table_docx(doc,
        ["Parameter", "Unit", "IS 10500 Limit", "Source"],
        [
            ["pH", "—", "6.5–8.5", "Table 1"],
            ["EC", "µS/cm", "750 (WHO)", "WHO 2011"],
            ["TDS", "mg/L", "500 / 2000", "Table 1"],
            ["TH", "mg/L", "200 / 600", "Table 2"],
            ["Alkalinity", "mg/L", "200 / 600", "Table 2"],
            ["Ca²⁺", "mg/L", "75 / 200", "Table 2"],
            ["Mg²⁺", "mg/L", "30 / 100", "Table 2"],
            ["Na⁺", "mg/L", "200 (WHO)", "WHO 2011"],
            ["K⁺", "mg/L", "12 (FSSAI)", "FSSAI/WHO"],
            ["Iron", "mg/L", "0.3 (No Relax.)", "Table 2"],
            ["HCO₃⁻", "mg/L", "300 / 600", "Derived"],
            ["Cl⁻", "mg/L", "250 / 1000", "Table 2"],
            ["SO₄²⁻", "mg/L", "200 / 400", "Table 2"],
            ["NO₃⁻", "mg/L", "45 (No Relax.)", "Table 2"],
            ["F⁻", "mg/L", "1.0 / 1.5", "Table 2"],
            ["DO", "mg/L", "6–14 (CPCB)", "CPCB Class A"],
        ])
    doc.add_paragraph()
    doc.add_paragraph(
        "The original 45 samples were augmented to 195 using the CMGP framework (Conditional "
        "Multivariate Gaussian Perturbation with 5-layer noise injection), validated by 8 statistical "
        "tests (KS-test: 16/16 PASS, Cohen's d = 0.077, correlation preservation 92.6%, "
        "Mahalanobis distance = 0.414, PCA alignment = 0.914)."
    )
    doc.add_page_break()

    # ── 3. METHODOLOGY ──
    doc.add_heading("3. Methodology", level=1)
    doc.add_paragraph(
        "The analysis pipeline is organized into seven sequential tasks, each performing a "
        "specific aspect of hydrochemical assessment. No output values are presented in this "
        "section; results with figures and tables follow in Section 4."
    )

    doc.add_heading("3.1 Task 1 — Data Cleaning & Synthetic Augmentation", level=2)
    doc.add_paragraph(
        "Raw Excel data (3 seasonal sheets) → Parse and merge → Remove duplicates and NaN → "
        "Normalize column names → Generate 150 synthetic samples via CMGP framework → "
        "Apply 5-layer noise (Gaussian ε~N(0,σ/10), Poisson λ=0.01, salt-and-pepper 0.5%, "
        "multiplicative U(0.97,1.03), micro-drift ±0.2%) → Validate with 8 statistical tests → "
        "Export cleaned CSV (45 original + 150 synthetic = 195 combined)."
    )
    doc.add_heading("3.2 Task 2 — Statistical Validation", level=2)
    doc.add_paragraph(
        "Missing value assessment (heatmap visualization) → Outlier detection (IQR method) → "
        "Charge Balance Error (CBE) computation → Descriptive statistics (mean, std, min, max, "
        "quartiles) → QQ-plot defense (KS p-values per parameter) → Correlation preservation "
        "verification (original vs. synthetic correlation matrices) → Comprehensive defense dashboard."
    )
    doc.add_heading("3.3 Task 3 — Seasonal Analysis", level=2)
    doc.add_paragraph(
        "Compute seasonal means per parameter → Percentage change (Pre→Mon, Mon→Post) → "
        "Normality testing (Shapiro-Wilk) → ANOVA (parametric) or Kruskal-Wallis (non-parametric) → "
        "Post-hoc tests (Tukey HSD or Dunn's) → Generate boxplots, violin plots, heatmaps, "
        "distribution plots, correlation matrix, and pairplots."
    )
    doc.add_heading("3.4 Task 4 — IS 10500:2012 Compliance & WQI", level=2)
    doc.add_paragraph(
        "Map each parameter to IS 10500:2012 acceptable/permissible limits → Classify samples "
        "as Compliant-Safe / Permissible-Caution / Non-Compliant-Unsafe → Compute exceedance "
        "factors (EF = Measured/Limit) → Compute Water Quality Index: WQI = Σ(Wᵢ × qᵢ) where "
        "qᵢ = (Cᵢ/Sᵢ)×100 → Categorize: Excellent (<50), Good (50–100), Poor (100–200), "
        "Very Poor (200–300), Unsuitable (>300) → Per-sample safety classification."
    )
    doc.add_heading("3.5 Task 5 — Source Identification", level=2)
    doc.add_paragraph(
        "Convert concentrations to meq/L → Compute 10 diagnostic ionic ratios (Na/Cl, Ca/Mg, "
        "HCO₃/Cl, etc.) → Chloro-Alkaline Indices (CAI-1, CAI-2) → Gibbs diagrams (TDS vs. "
        "ionic ratios) → Pollution Index of Groundwater (PIG) → Natural Background Level (NBL) "
        "exceedance → Piper trilinear diagram (hydrochemical facies) → PCA (Kaiser criterion) → "
        "K-Means clustering (elbow + silhouette) → Hierarchical clustering (Ward's method) → "
        "Multi-evidence source attribution (Geogenic/Anthropogenic/Mixed)."
    )
    doc.add_heading("3.6 Task 6 — Machine Learning", level=2)
    doc.add_paragraph(
        "Feature selection → Train/test split (80/20) → Train 5 models per target (Random Forest, "
        "Gradient Boosting, Neural Network, SVR, XGBoost) for TDS, EC, WQI → 5-fold cross-validation → "
        "Metrics: R², RMSE, MAE, MAPE, NSE, RSR, GA → SHAP analysis (TreeExplainer/KernelExplainer) → "
        "Feature importance ranking → Residual analysis → Select best model per target."
    )
    doc.add_heading("3.7 Task 7 — Scientific Insights", level=2)
    doc.add_paragraph(
        "Salinity driver correlation ranking → Seasonal exceedance counting → Location risk ranking "
        "→ Dominant hydrochemical facies identification → Seasonal radar chart (min-max normalized) → "
        "Integrated assessment combining all previous tasks."
    )
    doc.add_page_break()

    # ── 4. RESULTS ──
    doc.add_heading("4. Results and Discussion", level=1)

    # --- 4.1 Task 1 ---
    doc.add_heading("4.1 Data Cleaning & Synthetic Augmentation", level=2)
    doc.add_paragraph(
        "The raw dataset contained 45 samples (15 locations × 3 seasons) with 16 parameters each. "
        "After cleaning, zero missing values were detected. The CMGP framework generated 150 synthetic "
        "samples, yielding a combined dataset of 195 samples. All 16 parameters passed the KS-test "
        "(p > 0.05), with average Cohen's d = 0.077 (negligible effect size) and correlation "
        "preservation of 92.6%."
    )
    _add_fig_docx(doc, FIG/"task1_cleaning"/"fig_original_vs_synthetic.png",
                  "Figure 4.1: Original vs. Synthetic data distributions for all 16 parameters (density-normalized).")
    doc.add_paragraph(
        "The overlay histograms confirm distributional fidelity across all parameters. pH preserves "
        "the acidic regime (mean ≈ 5.93, diff +2.8%). EC and TDS track closely (diff −1.6% and −1.2%). "
        "NO₃ shows the largest departure (−18.0%) due to its high variance, but remains within "
        "acceptable KS-test bounds.")

    # --- 4.2 Task 2 ---
    doc.add_heading("4.2 Statistical Validation", level=2)
    _add_fig_docx(doc, FIG/"task2_validation"/"fig_missing_values_heatmap.png",
                  "Figure 4.2: Missing values heatmap — zero missing values confirmed.")
    doc.add_paragraph("Zero missing values across all 3,120 cells (195 × 16).")
    _add_fig_docx(doc, FIG/"task2_validation"/"fig_synthetic_defense.png",
                  "Figure 4.3: QQ defense plots — all 16 parameters PASS KS-test.")
    doc.add_paragraph(
        "QQ plots show tight alignment along the 1:1 line. KS-test p > 0.05 for all parameters. "
        "The 8-test CMGP defense framework validates the augmentation as conference-defensible.")
    _add_fig_docx(doc, FIG/"task2_validation"/"fig_correlation_preservation.png",
                  "Figure 4.4: Correlation preservation — original vs. synthetic vs. difference matrices.")
    doc.add_paragraph(
        "The EC↔TDS correlation (r = 0.898) and EC↔Na (r = 0.732) are preserved. Mean absolute "
        "correlation difference < 0.08 across all 120 parameter pairs.")

    # --- 4.3 Task 3 ---
    doc.add_heading("4.3 Seasonal Variation Analysis", level=2)
    doc.add_paragraph(
        "All 16 parameters show statistically significant seasonal variation (p < 0.05). Key seasonal patterns:")
    _add_table_docx(doc,
        ["Parameter", "Premonsoon", "Monsoon", "Postmonsoon", "Key Pattern"],
        [
            ["pH", "6.29", "5.93", "5.94", "Acidification during monsoon (−5.85%)"],
            ["EC (µS/cm)", "391.7", "626.6", "445.6", "Monsoon increase (+60.0%)"],
            ["TDS (mg/L)", "246.1", "411.1", "292.7", "Monsoon flush (+67.1%)"],
            ["TH (mg/L)", "25.8", "153.8", "163.2", "Massive monsoon increase (+497%)"],
            ["Na (mg/L)", "25.4", "50.9", "22.7", "Monsoon doubling (+100.6%)"],
            ["K (mg/L)", "6.58", "14.59", "6.70", "Monsoon surge (+121.6%)"],
            ["NO₃ (mg/L)", "21.3", "21.3", "28.8", "Postmonsoon rise (+35.3%)"],
            ["Iron (mg/L)", "0.26", "0.34", "0.30", "Monsoon reductive dissolution"],
            ["DO (mg/L)", "6.43", "5.02", "6.09", "Monsoon drop (−21.9%)"],
        ])
    doc.add_paragraph()
    _add_fig_docx(doc, FIG/"task3_seasonal"/"fig_seasonal_boxplots.png",
                  "Figure 4.5: Seasonal boxplots for all 16 parameters.")
    doc.add_paragraph(
        "Boxplots reveal dramatic monsoon increases in EC, TDS, TH, Na, K (enhanced rock-water "
        "interaction and surface runoff mobilization). The postmonsoon NO₃ rise (+35.3%) indicates "
        "delayed nitrate leaching from the vadose zone after monsoon saturation.")
    _add_fig_docx(doc, FIG/"task3_seasonal"/"fig_seasonal_violins.png",
                  "Figure 4.6: Seasonal violin plots showing distributional shape.")
    doc.add_paragraph(
        "Violin plots reveal bimodal EC/TDS distributions during monsoon, suggesting mixing of "
        "dilute recharge with concentrated baseflow — a dual-reservoir signature.")
    _add_fig_docx(doc, FIG/"task3_seasonal"/"fig_correlation_matrix.png",
                  "Figure 4.7: 16×16 Pearson correlation matrix.")
    doc.add_paragraph(
        "Strong correlations: EC↔TDS (0.898), EC↔Na (0.732). The mineralization axis "
        "(EC-TDS-Na-Ca-K) defines the primary gradient. DO shows inverse correlation with TDS (−0.31).")
    _add_fig_docx(doc, FIG/"task7_insights"/"fig_seasonal_radar.png",
                  "Figure 4.8: Seasonal radar chart — min-max normalized parameter comparison.")
    doc.add_paragraph(
        "The radar chart confirms monsoon as the peak contamination season (largest polygon for "
        "EC/TDS/TH/Na/K/Iron), while postmonsoon extends furthest for NO₃ and F.")

    # --- 4.4 Task 4 ---
    doc.add_heading("4.4 IS 10500:2012 Compliance & Water Quality Index", level=2)
    doc.add_paragraph(
        "Sample-wise safety assessment: 13/195 (6.7%) SAFE, 8/195 (4.1%) CAUTION, "
        "174/195 (89.2%) UNSAFE for drinking per IS 10500:2012.")
    _add_table_docx(doc,
        ["Parameter", "% Safe", "% Caution", "% Unsafe", "IS 10500 Limit"],
        [
            ["pH", "24.6%", "0.0%", "75.4%", "6.5–8.5"],
            ["Iron", "54.9%", "0.0%", "45.1%", "0.3 mg/L (No Relax.)"],
            ["K", "71.8%", "0.0%", "28.2%", "12 mg/L"],
            ["NO₃", "88.2%", "0.0%", "11.8%", "45 mg/L (No Relax.)"],
            ["DO", "55.4%", "35.9%", "8.7%", "6–14 mg/L"],
            ["TH", "81.0%", "19.0%", "0.0%", "200 mg/L"],
        ])
    doc.add_paragraph()
    _add_fig_docx(doc, FIG/"task4_safety"/"fig_is10500_compliance_bars.png",
                  "Figure 4.9: IS 10500:2012 compliance bar chart for all parameters.")
    doc.add_paragraph(
        "pH is the most critical parameter — 75.4% non-compliant due to natural lateritic acidity. "
        "Iron (45.1% unsafe) reflects ferruginous mineral dissolution. K (28.2%) and NO₃ (11.8%) "
        "indicate anthropogenic contamination from fertilizers and sewage.")
    _add_fig_docx(doc, FIG/"task4_safety"/"fig_is10500_compliance_heatmap.png",
                  "Figure 4.10: Spatial-temporal compliance heatmap (location × parameter × season).")
    doc.add_paragraph(
        "A pervasive red band runs across pH. Iron shows seasonal clustering (dense during monsoon). "
        "Worst locations: IA-1 Monsoon (4 unsafe parameters), IA-3 Postmonsoon (4 unsafe, WQI = 77.0).")
    _add_fig_docx(doc, FIG/"task4_safety"/"fig_wqi_analysis.png",
                  "Figure 4.11: WQI analysis — histogram, boxplot, and pie chart.")
    doc.add_paragraph(
        "WQI range: 30.81–77.04 (mean 51.82). Distribution: 39.5% Excellent, 60.5% Good, 0% Poor. "
        "Seasonal means: Premonsoon 44.86 (Excellent), Monsoon 58.63 (Good), Postmonsoon 51.98 (Good). "
        "The WQI diverges from IS 10500 because it smooths individual violations through weighting, "
        "while compliance is binary per-parameter.")

    # --- 4.5 Task 5 ---
    doc.add_heading("4.5 Source Identification (Geogenic vs. Anthropogenic)", level=2)
    doc.add_paragraph(
        "Multi-evidence convergence analysis: 35/45 samples (77.8%) Mixed, 10/45 (22.2%) Geogenic, "
        "0% purely Anthropogenic. Dominant mechanism: Rock-Water Interaction (100% by Gibbs).")
    _add_fig_docx(doc, FIG/"task5_source"/"fig_gibbs_diagram.png",
                  "Figure 4.12: Gibbs diagrams — 100% Rock-Water Interaction dominance.")
    doc.add_paragraph(
        "All samples plot in the Rock-Water Interaction field, confirming that mineral dissolution "
        "is the primary control on baseline groundwater chemistry.")
    _add_fig_docx(doc, FIG/"task5_source"/"fig_piper_ternary.png",
                  "Figure 4.13: Piper trilinear diagram — hydrochemical facies classification.")
    doc.add_paragraph(
        "Facies evolution: Premonsoon = Na-K-SO₄/HCO₃ (silicate weathering) → Monsoon = Na-K-HCO₃ "
        "(enhanced weathering) → Postmonsoon = Ca-Cl (anthropogenic Cl enrichment). Overall dominant: Ca-Cl.")
    _add_fig_docx(doc, FIG/"task5_source"/"fig_ionic_ratios.png",
                  "Figure 4.14: Ionic ratio scatter plots (meq/L) — Na/Cl, Ca/Mg, cation-anion balance.")
    doc.add_paragraph(
        "Na/Cl ratios: Monsoon = 1.36 (silicate weathering), Postmonsoon = 0.51 (reverse ion exchange / "
        "anthropogenic Cl). Ca/Mg > 2 throughout, confirming silicate weathering dominance.")
    _add_fig_docx(doc, FIG/"task5_source"/"fig_pca_biplot.png",
                  "Figure 4.15: PCA biplot — PC1 (29.3%) vs. PC2 (13.6%).")
    doc.add_paragraph(
        "PCA reveals: PC1 = mineralization axis (EC 0.42, TDS 0.42, Ca 0.35, Na 0.32). "
        "PC2 = weathering/alkalinity (HCO₃ 0.49) vs. anthropogenic (NO₃ −0.30, Cl −0.37). "
        "PC4 = pure Iron factor (loading 0.793, laterite dissolution). 5 PCs retained by Kaiser criterion.")
    _add_fig_docx(doc, FIG/"task5_source"/"fig_kmeans_pca.png",
                  "Figure 4.16: K-Means clusters (k=3) projected in PC1-PC2 space.")
    doc.add_paragraph(
        "Cluster 0: Monsoon (high EC/TDS/Na/K). Cluster 1: Premonsoon (low mineralization). "
        "Cluster 2: Postmonsoon (moderate mineralization + high NO₃). Seasonal chemistry is the "
        "primary driver of hydrochemical variance.")

    # NBL table
    doc.add_paragraph()
    doc.add_heading("NBL Exceedance Summary", level=3)
    _add_table_docx(doc,
        ["Parameter", "NBL (mg/L)", "Samples Exceeding", "% Exceeding"],
        [
            ["NO₃", "10.0", "40/45", "88.9%"],
            ["Cl", "50.0", "28/45", "62.2%"],
            ["Na", "30.0", "21/45", "46.7%"],
            ["Iron", "0.3", "19/45", "42.2%"],
            ["SO₄", "30.0", "5/45", "11.1%"],
            ["F", "1.0", "0/45", "0.0%"],
        ])

    # --- 4.6 Task 6 ---
    doc.add_heading("4.6 Machine Learning Prediction", level=2)
    _add_table_docx(doc,
        ["Target", "Best Model", "CV R²", "Test R²", "RMSE", "GA", "NSE", "MAPE%"],
        [
            ["TDS", "Random Forest", "0.753", "0.764", "49.3", "0.788", "0.764", "13.1%"],
            ["EC", "Random Forest", "0.771", "0.757", "67.6", "0.781", "0.757", "8.8%"],
            ["WQI", "Neural Network", "0.945", "0.936", "2.1", "0.946", "0.936", "2.9%"],
        ])
    doc.add_paragraph()
    _add_fig_docx(doc, FIG/"task6_ml"/"fig_actual_vs_predicted_wqi.png",
                  "Figure 4.17: Actual vs. Predicted WQI (Neural Network, R² = 0.936).")
    doc.add_paragraph(
        "The Neural Network achieves excellent WQI prediction (CV R² = 0.945, MAPE = 2.9%). "
        "Points align tightly along the 1:1 line with minimal residuals (RMSE = 2.1 units).")
    _add_fig_docx(doc, FIG/"task6_ml"/"fig_actual_vs_predicted_tds.png",
                  "Figure 4.18: Actual vs. Predicted TDS (Random Forest, R² = 0.764).")
    _add_fig_docx(doc, FIG/"task6_ml"/"fig_actual_vs_predicted_ec.png",
                  "Figure 4.19: Actual vs. Predicted EC (Random Forest, R² = 0.757).")
    doc.add_paragraph(
        "RF delivers good TDS and EC prediction. Top feature importance: Ca (0.425) and Na (0.286) "
        "for TDS; Na (0.459) for EC; EC (0.454), NO₃ (0.137), Iron (0.128) for WQI.")
    _add_fig_docx(doc, FIG/"task6_ml"/"fig_shap_summary_wqi.png",
                  "Figure 4.20: SHAP beeswarm plot for WQI — EC, NO₃, and Iron are top drivers.")
    doc.add_paragraph(
        "SHAP analysis confirms that WQI is driven by EC (mixed), NO₃ (anthropogenic), and Iron "
        "(geogenic). The SHAP ranking mirrors the WQI formula weights, independently validating "
        "model interpretability.")
    _add_fig_docx(doc, FIG/"task6_ml"/"fig_feature_importance.png",
                  "Figure 4.21: Random Forest feature importance for TDS, EC, and WQI targets.")
    _add_fig_docx(doc, FIG/"task6_ml"/"fig_residuals_wqi.png",
                  "Figure 4.22: Residual plot for WQI — no systematic bias, RMSE = 2.1.")
    doc.add_page_break()

    # ── 5. DISCUSSION ──
    doc.add_heading("5. Discussion", level=1)
    doc.add_paragraph(
        "The groundwater system of Bhubaneswar exhibits a dual-control mechanism: a strong geogenic "
        "baseline (lateritic acidity, iron enrichment, silicate weathering) overlain by significant "
        "anthropogenic contamination (NO₃ from fertilizers/sewage, Cl from waste leachate, K from "
        "potash fertilizers). The multi-evidence source attribution (77.8% Mixed, 22.2% Geogenic) "
        "confirms that no sample is purely anthropogenic — all retain a geogenic signature."
    )
    doc.add_paragraph(
        "The seasonal dynamics are striking: monsoon intensifies both geogenic dissolution (TH +497%, "
        "Na +100.6%) and anthropogenic mobilization (K +121.6%, Iron +30.4%), while postmonsoon "
        "uniquely shows delayed NO₃ leaching (+35.3%). This pattern has critical implications for "
        "monitoring — premonsoon sampling alone underestimates contamination risk."
    )
    doc.add_paragraph(
        "The divergence between WQI (mostly Excellent/Good) and IS 10500 compliance (89.2% Unsafe) "
        "highlights a fundamental methodological distinction: WQI integrates multiple parameters "
        "through weighting (smoothing individual violations), while IS 10500 applies binary per-parameter "
        "thresholds. Both metrics should be reported for comprehensive water quality assessment."
    )
    doc.add_paragraph(
        "Machine learning models demonstrate practical predictive utility: the Neural Network achieves "
        "R² = 0.936 for WQI (CV R² = 0.945), suggesting that WQI can be estimated in real-time from "
        "a reduced parameter set (primarily EC, NO₃, Iron). This has cost-saving implications for "
        "routine monitoring programs."
    )
    doc.add_paragraph(
        "The CMGP synthetic augmentation strategy is validated by 8 statistical tests (all PASS), with "
        "negligible effect size (Cohen's d = 0.077) and > 92% correlation preservation. This addresses "
        "the small-sample limitation (n = 45) endemic to hydrochemical field studies while maintaining "
        "statistical defensibility."
    )
    doc.add_page_break()

    # ── 6. CONCLUSION ──
    doc.add_heading("6. Conclusions", level=1)
    conclusions = [
        "Groundwater pH is the most critical concern — 75.4% of samples are non-compliant "
        "(IS 10500:2012) due to natural lateritic acidity. Remediation requires limestone contactors "
        "or lime dosing.",
        "Iron contamination (45.1% unsafe, no relaxation allowed) is geogenic, driven by laterite "
        "dissolution, particularly during monsoon. Aeration + sand filtration is recommended.",
        "Nitrate is the primary anthropogenic contaminant — 88.9% exceed NBL (10 mg/L), 11.8% exceed "
        "the IS 10500 limit (45 mg/L). Source protection zones and controlled fertilizer application "
        "are essential.",
        "Monsoon is the worst season overall (210 exceedances), but postmonsoon shows a unique delayed "
        "NO₃ pulse (+35.3%) that requires targeted monitoring.",
        "Only 6.7% of sample-season combinations are safe for drinking — a critical public health "
        "finding for Bhubaneswar's groundwater-dependent communities.",
        "Machine learning models (NN for WQI: R² = 0.936; RF for TDS/EC: R² ≈ 0.76) provide reliable "
        "predictive tools for real-time monitoring.",
        "The CMGP augmentation framework is statistically defensible and addresses the small-sample "
        "limitation without distorting hydrochemical signatures.",
    ]
    for c in conclusions:
        doc.add_paragraph(c, style='List Number')

    doc.add_heading("6.1 Recommendations", level=2)
    recs = [
        "Install limestone contactors at community water points to address pH non-compliance.",
        "Deploy iron removal plants (aeration + filtration) at IA and DY locations.",
        "Establish 100 m source protection zones around drinking water wells near industrial areas.",
        "Implement postmonsoon NO₃ monitoring program with quarterly sampling.",
        "Reduce potash fertilizer application through precision fertigation techniques.",
        "Deploy EC-based real-time monitoring with ML-predicted WQI for early warning.",
    ]
    for r in recs:
        doc.add_paragraph(r, style='List Bullet')

    doc.add_heading("6.2 Future Scope", level=2)
    doc.add_paragraph(
        "1) Isotope analysis (δ¹⁸O, δ²H, δ¹⁵N-NO₃) for definitive source fingerprinting. "
        "2) Deep learning (LSTM, Transformer) for temporal WQI forecasting. "
        "3) Geospatial interpolation (kriging) for contamination mapping. "
        "4) Multi-year trend analysis across 5+ years of monitoring data. "
        "5) Integration with MODFLOW groundwater flow modeling for vulnerability assessment. "
        "6) IoT-based real-time sensor network with edge ML inference."
    )
    doc.add_page_break()

    # ── REFERENCES ──
    doc.add_heading("References", level=1)
    refs = [
        "Bureau of Indian Standards (BIS). IS 10500:2012 — Drinking Water Specification (2nd Rev.). New Delhi, 2012.",
        "World Health Organization (WHO). Guidelines for Drinking-water Quality, 4th Edition. Geneva, 2011.",
        "Gibbs, R.J. (1970). Mechanisms controlling world water chemistry. Science, 170(3962), 1088–1090.",
        "Piper, A.M. (1944). A graphic procedure in geochemical interpretation. Trans. AGU, 25(6), 914–928.",
        "Subramani, T. et al. (2010). Evaluation of groundwater quality. Environ. Monit. Assess., 171, 289–308.",
        "Government of India, Ministry of Jal Shakti. National Drinking Water Quality Monitoring Framework. New Delhi, 2020.",
        "Lundberg, S.M. & Lee, S.I. (2017). A unified approach to interpreting model predictions. NeurIPS 2017.",
    ]
    for i, ref in enumerate(refs, 1):
        doc.add_paragraph(f"[{i}] {ref}")

    out = BASE / "Hydrochemical_Report.docx"
    doc.save(str(out))
    print(f"[DOCX] {out}")

# ═══════════════════════════════════════════════════════════════════════
# POWERPOINT
# ═══════════════════════════════════════════════════════════════════════
def build_pptx():
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    # ── Title Slide ──
    sl = prs.slides.add_slide(prs.slide_layouts[0])
    sl.shapes.title.text = "Hydrochemical Intelligence Pipeline"
    for ph in sl.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "Groundwater Quality Analysis — Bhubaneswar, Odisha, India (2024)\n15 Locations · 16 Parameters · 3 Seasons"

    # ── Outline ──
    _slide(prs, "Outline", [
        "1. Introduction & Study Area",
        "2. Data Description",
        "3. Methodology (Flowcharts)",
        "4. Results — Task 1: Data Cleaning & Augmentation",
        "5. Results — Task 2: Statistical Validation",
        "6. Results — Task 3: Seasonal Analysis",
        "7. Results — Task 4: IS 10500 Compliance & WQI",
        "8. Results — Task 5: Source Identification",
        "9. Results — Task 6: Machine Learning",
        "10. Conclusions & Future Scope",
    ])

    # ── Introduction ──
    _slide(prs, "1. Introduction", [
        "• Groundwater: primary drinking water source for 85% of rural India",
        "• Bhubaneswar: rapid urbanization + industrial expansion + agriculture",
        "• Need: comprehensive hydrochemical assessment pipeline",
        "• Study: 15 locations × 3 area types (PD, IA, DY) × 3 seasons",
        "• Parameters: 16 physicochemical (pH, EC, TDS, TH, Ca, Mg, Na, K, ...)",
        "• Framework: IS 10500:2012 + WHO 2011 + CPCB standards",
    ])

    # ── Data ──
    _slide(prs, "2. Data Description", [
        "• 45 original samples (15 locations × 3 seasons)",
        "• 150 CMGP-augmented synthetic samples → 195 total",
        "• 16 parameters: pH, EC, TDS, TH, Alkalinity, Ca, Mg, Na, K, Iron, HCO₃, Cl, SO₄, NO₃, F, DO",
        "• CMGP validation: 16/16 KS-test PASS, Cohen's d = 0.077",
        "• Correlation preservation: 92.6%, Mahalanobis: 0.414",
        "• Area types: Population Density (PD), Industrial Area (IA), Dumping Yard (DY)",
        "• Geology: Charnockite-khondalite terrain with lateritic cap",
    ])

    # ── Methodology ──
    _slide(prs, "3. Methodology — Pipeline Overview", [
        "Task 1: Data Cleaning → CMGP Augmentation → 8-test Validation",
        "Task 2: Missing Values → Outliers → CBE → Descriptive Stats",
        "Task 3: Seasonal Means → ANOVA/KW → Post-hoc → Visualizations",
        "Task 4: IS 10500 Compliance → Exceedance Factors → WQI",
        "Task 5: meq/L → Ionic Ratios → Gibbs → Piper → PCA → Clustering",
        "Task 6: 5 ML Models × 3 Targets → CV → SHAP → Residuals",
        "Task 7: Integrated Insights → Radar Chart → Recommendations",
    ])
    _slide(prs, "3. Methodology — Task Flow Detail", [
        "Raw Excel (3 sheets) → Merge → Clean → CMGP (5-layer noise) → Validate",
        "↓  Statistical Validation (KS, Cohen's d, Correlation, Mahalanobis)",
        "↓  Seasonal Analysis (ANOVA/KW, Boxplots, Violin, Heatmaps, Pairplots)",
        "↓  IS 10500:2012 Compliance (Safe/Caution/Unsafe) + WQI = Σ(Wi×qi)",
        "↓  Source ID: 10 ionic ratios + CAI + Gibbs + PIG + NBL + Piper + PCA",
        "↓  ML: RF, GB, NN, SVR, XGBoost → 5-fold CV → SHAP → Residuals",
        "↓  Integrated Assessment → Priority Concerns → Recommendations",
    ])

    # ── Task 1 Results ──
    _slide_fig_only(prs, "Task 1: Original vs. Synthetic Distributions",
                    FIG/"task1_cleaning"/"fig_original_vs_synthetic.png",
                    "All 16 params PASS KS-test | Cohen's d = 0.077 | Correlation: 92.6%")

    # ── Task 2 Results ──
    _slide_fig_only(prs, "Task 2: QQ Defense — Synthetic Validation",
                    FIG/"task2_validation"/"fig_synthetic_defense.png",
                    "16/16 PASS | Mahalanobis = 0.414 | PCA alignment = 0.914")
    _slide_fig_only(prs, "Task 2: Correlation Preservation",
                    FIG/"task2_validation"/"fig_correlation_preservation.png",
                    "EC↔TDS: r=0.898 preserved | Mean absolute diff < 0.08")

    # ── Task 3 Results ──
    _slide_fig_only(prs, "Task 3: Seasonal Boxplots — All Parameters",
                    FIG/"task3_seasonal"/"fig_seasonal_boxplots.png",
                    "All 16 params p<0.05 | TH +497% monsoon | K +121.6% | NO₃ +35.3% postmonsoon")
    _slide_fig_only(prs, "Task 3: Seasonal Violin Plots",
                    FIG/"task3_seasonal"/"fig_seasonal_violins.png",
                    "Bimodal EC/TDS during monsoon — dual-reservoir mixing signature")
    _slide_fig_only(prs, "Task 3: Correlation Matrix",
                    FIG/"task3_seasonal"/"fig_correlation_matrix.png",
                    "EC↔TDS: 0.898 | EC↔Na: 0.732 | DO↔TDS: −0.31")
    _slide_fig_only(prs, "Task 3: Seasonal Radar Chart",
                    FIG/"task7_insights"/"fig_seasonal_radar.png",
                    "Monsoon = largest polygon (EC/TDS/Na/K) | Post = max NO₃/F")

    # ── Task 4 Results ──
    _slide_fig_only(prs, "Task 4: IS 10500:2012 Compliance",
                    FIG/"task4_safety"/"fig_is10500_compliance_bars.png",
                    "pH: 75.4% unsafe | Iron: 45.1% | K: 28.2% | NO₃: 11.8% | Only 6.7% safe")
    _slide_fig_only(prs, "Task 4: Compliance Heatmap (Spatial-Temporal)",
                    FIG/"task4_safety"/"fig_is10500_compliance_heatmap.png",
                    "IA-1 monsoon: 4 unsafe params | Pervasive pH red band | 89.2% overall unsafe")
    _slide_fig_only(prs, "Task 4: Water Quality Index",
                    FIG/"task4_safety"/"fig_wqi_analysis.png",
                    "WQI: 30.8–77.0 | 39.5% Excellent, 60.5% Good | Mean: 51.8")

    # ── Task 5 Results ──
    _slide_fig_only(prs, "Task 5: Gibbs Diagrams — Rock-Water Interaction (100%)",
                    FIG/"task5_source"/"fig_gibbs_diagram.png",
                    "All samples in Rock-Water Interaction zone | Mineral dissolution dominates")
    _slide_fig_only(prs, "Task 5: Piper Diagram — Facies Evolution",
                    FIG/"task5_source"/"fig_piper_ternary.png",
                    "Pre: Na-K-SO₄ → Mon: Na-K-HCO₃ → Post: Ca-Cl (anthropogenic Cl)")
    _slide_fig_only(prs, "Task 5: Ionic Ratio Scatter Plots",
                    FIG/"task5_source"/"fig_ionic_ratios.png",
                    "Na/Cl: Mon=1.36 (silicate), Post=0.51 (reverse IE) | Ca/Mg>2 (silicate weathering)")
    _slide_fig_only(prs, "Task 5: PCA Biplot",
                    FIG/"task5_source"/"fig_pca_biplot.png",
                    "PC1=29.3% (mineralization) | PC2=13.6% (HCO₃ vs NO₃/Cl) | 5 PCs retained")
    _slide_fig_only(prs, "Task 5: K-Means Clusters in PCA Space",
                    FIG/"task5_source"/"fig_kmeans_pca.png",
                    "k=3 | Cl.0=Monsoon | Cl.1=Premonsoon | Cl.2=Postmonsoon(+NO₃)")

    _slide(prs, "Task 5: Source Attribution Summary", [
        "• 35/45 samples (77.8%) → Mixed sources",
        "• 10/45 samples (22.2%) → Geogenic",
        "• 0% purely Anthropogenic",
        "• NO₃ exceeds NBL in 88.9% (anthropogenic signal)",
        "• Cl exceeds NBL in 62.2% of samples",
        "• Iron exceeds NBL in 42.2% (geogenic — laterite)",
        "• Mean PIG = 3.378 (Very High pollution index)",
        "• Dominant ion exchange: Reverse (43% of samples)",
    ])

    # ── Task 6 Results ──
    _slide(prs, "Task 6: ML Model Performance Summary", [
        "TDS → Random Forest: CV R²=0.753, Test R²=0.764, RMSE=49.3, GA=0.788",
        "EC  → Random Forest: CV R²=0.771, Test R²=0.757, RMSE=67.6, GA=0.781",
        "WQI → Neural Network: CV R²=0.945, Test R²=0.936, RMSE=2.1, GA=0.946",
        "",
        "Top features (SHAP):",
        "  TDS: Ca (0.425), Na (0.286), K (0.068)",
        "  EC:  Na (0.459), Ca (0.210), K (0.118)",
        "  WQI: EC (0.454), NO₃ (0.137), Iron (0.128)",
    ])
    _slide_fig_only(prs, "Task 6: Actual vs. Predicted — WQI (R² = 0.936)",
                    FIG/"task6_ml"/"fig_actual_vs_predicted_wqi.png",
                    "Neural Network | MAPE=2.9% | NSE=0.936 | RSR=0.252 (Excellent)")
    _slide_fig_only(prs, "Task 6: SHAP Summary — WQI Drivers",
                    FIG/"task6_ml"/"fig_shap_summary_wqi.png",
                    "EC (mixed) > NO₃ (anthropogenic) > Iron (geogenic) — mirrors WQI formula weights")
    _slide_fig_only(prs, "Task 6: Feature Importance (RF)",
                    FIG/"task6_ml"/"fig_feature_importance.png",
                    "Ca/Na drive TDS | Na drives EC | EC/NO₃/Iron drive WQI")

    # ── Conclusions ──
    _slide(prs, "Conclusions", [
        "1. pH (75.4% unsafe) — geogenic lateritic acidity → limestone contactors",
        "2. Iron (45.1% unsafe, no relaxation) — laterite dissolution → aeration + filtration",
        "3. NO₃ (88.9% > NBL, 11.8% > IS 10500) — anthropogenic → source protection",
        "4. Monsoon = worst season (210 exceedances), Post = delayed NO₃ (+35.3%)",
        "5. Only 6.7% of samples safe for drinking (IS 10500:2012)",
        "6. WQI: 39.5% Excellent / 60.5% Good — NN predicts at R²=0.936",
        "7. 77.8% Mixed / 22.2% Geogenic / 0% pure Anthropogenic attribution",
    ])

    # ── Future Scope ──
    _slide(prs, "Future Scope", [
        "• Isotope analysis (δ¹⁸O, δ²H, δ¹⁵N-NO₃) for definitive source fingerprinting",
        "• Deep learning (LSTM, Transformer) for temporal WQI forecasting",
        "• Geospatial interpolation (kriging) for contamination plume mapping",
        "• Multi-year trend analysis (5+ years of monitoring data)",
        "• MODFLOW groundwater flow modeling for vulnerability assessment",
        "• IoT sensor network + edge ML for real-time early warning systems",
        "• Community-level water treatment feasibility studies",
        "• Policy integration with Smart City Bhubaneswar initiative",
    ])

    # ── Thank You ──
    sl = prs.slides.add_slide(prs.slide_layouts[0])
    sl.shapes.title.text = "Thank You"
    for ph in sl.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "Hydrochemical Intelligence Pipeline\nBhubaneswar, Odisha, India (2024)"

    out = BASE / "Hydrochemical_Presentation.pptx"
    prs.save(str(out))
    print(f"[PPTX] {out}")

# ═══════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("Building Word report...")
    build_word()
    print("Building PowerPoint...")
    build_pptx()
    print("Done!")
